# proposal_generator.py
import re
import os
from pptx import Presentation
from pptx.util import Pt
from io import BytesIO

def generate_proposal_ppt_data(generate_section_content, slide_section_mapping, firm, client_name, product, proposal_date, template_file=r"C:\Users\ssuji\OneDrive\Documents\GS-final\Proposal Template.pptx"):
    """
    Generate a proposal PPTX using the provided section content generator and slide mapping.
    Instead of collecting inputs via input(), use the passed parameters. 
    The generated presentation is returned as a BytesIO stream.
    """
    if os.path.exists(template_file):
        prs = Presentation(template_file)
    else:
        prs = Presentation()
    
    # Update title slide (assumed slide index 0)
    if prs.slides:
        title_slide = prs.slides[0]
        if title_slide.shapes.title:
            title_slide.shapes.title.text = f"{product} Proposal"
            for paragraph in title_slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(72)
        if len(title_slide.placeholders) > 1:
            subtitle_text = f"Consultancy Firm: {firm}\nClient: {client_name}\nProposal Date: {proposal_date}"
            title_slide.placeholders[1].text = subtitle_text
            for paragraph in title_slide.placeholders[1].text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(24)
    else:
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = f"{product} Proposal"
        for paragraph in title_slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(72)
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = f"Consultancy Firm: {firm}\nClient: {client_name}\nProposal Date: {proposal_date}"
            for paragraph in title_slide.placeholders[1].text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(24)
    
    max_bullets_per_slide = 8
    bullet_slide_layout = prs.slide_layouts[1]  # Assuming title and content layout
    
    for slide_idx, section in sorted(slide_section_mapping.items()):
        guideline = "Provide a concise summary."
        # Generate section content using the passed function from proposal_service
        section_content = generate_section_content(section, firm, client_name, product, proposal_date, guideline)
        bullet_points = re.split(r'\.\s+', section_content)
        bullet_points = [pt.strip() for pt in bullet_points if pt.strip()]
        
        # Use existing slide if available; otherwise add new slide
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            slide.shapes.title.text = section
            content_placeholder = slide.shapes.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
        else:
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = section
            content_placeholder = slide.shapes.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
        
        # Add bullet points; if too many bullets, create additional slides
        num_slides = (len(bullet_points) - 1) // max_bullets_per_slide + 1
        for i in range(num_slides):
            if i > 0:
                slide = prs.slides.add_slide(bullet_slide_layout)
                slide.shapes.title.text = f"{section} (Continued)"
                content_placeholder = slide.shapes.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
            start_idx = i * max_bullets_per_slide
            end_idx = start_idx + max_bullets_per_slide
            for point in bullet_points[start_idx:end_idx]:
                if not point.endswith('.'):
                    point += '.'
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(12)
    
    # Save the presentation into a BytesIO buffer
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

if __name__ == "__main__":
    # Dummy function for testing standalone
    def dummy_generate_section_content(section, firm, client_name, product, proposal_date, guideline):
        return f"This is dummy content for {section}."
    
    dummy_slide_section_mapping = {1: "Executive Summary", 2: "About Us"}
    # For standalone testing, you can uncomment and run:
    # generate_proposal_ppt_data(dummy_generate_section_content, dummy_slide_section_mapping, "Test Firm", "Test Client", "Test Product", "01-01-2025")
